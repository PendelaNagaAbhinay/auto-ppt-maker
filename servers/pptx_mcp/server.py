"""
PPTX MCP Server — Generates PowerPoint presentations via FastMCP tools.

Run:
    python server.py
"""

import json
import logging
import os
import traceback
from typing import Optional

from fastmcp import FastMCP
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ──────────────────────────────────────────────
# Logging  (writes to file — stderr is reserved for MCP stdio transport)
# ──────────────────────────────────────────────
_log_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "pptx_mcp.log")
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s | %(levelname)-7s | %(message)s",
    datefmt="%H:%M:%S",
    filename=_log_path,
    filemode="a",
)
logger = logging.getLogger("pptx_mcp")

# ──────────────────────────────────────────────
# Constants
# ──────────────────────────────────────────────
WIDESCREEN_WIDTH = Emu(12192000)   # 16:9 width  (13.333 in)
WIDESCREEN_HEIGHT = Emu(6858000)   # 16:9 height ( 7.5   in)

TITLE_FONT_SIZE = Pt(36)
BODY_FONT_SIZE = Pt(20)

MIN_BULLETS = 3
MAX_BULLETS = 5
MAX_BULLET_CHARS = 120

# Default theme colours
DEFAULT_BG = "FFFFFF"
DEFAULT_TITLE_COLOR = "1A1A2E"
DEFAULT_BODY_COLOR = "333333"

# ──────────────────────────────────────────────
# In-memory state
# ──────────────────────────────────────────────
_state: dict = {
    "prs": None,
    "slide_count": 0,
    "theme": {
        "background_color": DEFAULT_BG,
        "title_color": DEFAULT_TITLE_COLOR,
        "body_color": DEFAULT_BODY_COLOR,
    },
}

# ──────────────────────────────────────────────
# Helpers
# ──────────────────────────────────────────────

def _hex_to_rgb(hex_str: str) -> RGBColor:
    """Convert a hex colour string (with or without '#') to RGBColor."""
    hex_str = hex_str.lstrip("#")
    if len(hex_str) != 6:
        raise ValueError(f"Invalid hex colour: {hex_str}")
    r, g, b = int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16)
    return RGBColor(r, g, b)


def _require_presentation() -> Optional[str]:
    """Return an error string if no presentation exists, else None."""
    if _state["prs"] is None:
        return json.dumps({
            "status": "error",
            "message": "No presentation created yet. Call create_presentation() first.",
        })
    return None


def _sanitise_bullets(bullets: list[str]) -> tuple[list[str], list[str]]:
    """
    Enforce bullet rules:
      • 3–5 items  → pad with empty or trim
      • Each item ≤ 120 chars → truncate with '…'
    Returns (cleaned_bullets, warnings).
    """
    warnings: list[str] = []

    # Truncate individual bullets
    cleaned: list[str] = []
    for i, b in enumerate(bullets):
        if len(b) > MAX_BULLET_CHARS:
            trimmed = b[: MAX_BULLET_CHARS - 1] + "…"
            warnings.append(f"Bullet {i + 1} trimmed from {len(b)} to {MAX_BULLET_CHARS} chars.")
            cleaned.append(trimmed)
        else:
            cleaned.append(b)

    # Enforce count limits
    if len(cleaned) > MAX_BULLETS:
        warnings.append(f"Too many bullets ({len(cleaned)}). Trimmed to first {MAX_BULLETS}.")
        cleaned = cleaned[:MAX_BULLETS]
    elif len(cleaned) < MIN_BULLETS:
        warnings.append(f"Too few bullets ({len(cleaned)}). Padded to {MIN_BULLETS} with empty items.")
        while len(cleaned) < MIN_BULLETS:
            cleaned.append("")

    return cleaned, warnings


def _apply_background(slide) -> None:
    """Fill slide background with the current theme colour."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = _hex_to_rgb(_state["theme"]["background_color"])


def _add_slide_number(slide, slide_num: int) -> None:
    """Add a small slide-number footer in the bottom-right corner."""
    left = WIDESCREEN_WIDTH - Inches(1.2)
    top = WIDESCREEN_HEIGHT - Inches(0.5)
    txBox = slide.shapes.add_textbox(left, top, Inches(1.0), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = str(slide_num)
    p.alignment = PP_ALIGN.RIGHT
    run = p.runs[0]
    run.font.size = Pt(12)
    run.font.color.rgb = _hex_to_rgb(_state["theme"]["body_color"])


# ──────────────────────────────────────────────
# FastMCP Server
# ──────────────────────────────────────────────
mcp = FastMCP("pptx_mcp")


@mcp.tool
def create_presentation(theme: str = "default") -> str:
    """
    Create a new blank widescreen (16:9) presentation.

    Args:
        theme: Currently only 'default' is supported.

    Returns:
        JSON string confirming creation.
    """
    logger.info("create_presentation(theme=%s)", theme)

    prs = Presentation()
    prs.slide_width = WIDESCREEN_WIDTH
    prs.slide_height = WIDESCREEN_HEIGHT

    _state["prs"] = prs
    _state["slide_count"] = 0
    _state["theme"] = {
        "background_color": DEFAULT_BG,
        "title_color": DEFAULT_TITLE_COLOR,
        "body_color": DEFAULT_BODY_COLOR,
    }

    logger.info("Presentation created (16:9, theme=%s).", theme)
    return json.dumps({
        "status": "success",
        "message": f"Presentation created with theme '{theme}' (16:9 widescreen).",
        "slide_count": 0,
    })


@mcp.tool
def add_slide(title: str, bullets: list[str]) -> str:
    """
    Add a text slide with a title and 3-5 bullet points.

    Bullets are auto-trimmed if they exceed limits:
      - Max 5 items (extras dropped), min 3 (padded with blanks)
      - Each bullet max 120 characters (truncated with '…')

    Args:
        title:   Slide title text.
        bullets: List of bullet-point strings.

    Returns:
        JSON string with slide details and any warnings.
    """
    logger.info("add_slide(title=%r, bullets=%d items)", title, len(bullets))

    err = _require_presentation()
    if err:
        return err

    prs: Presentation = _state["prs"]
    cleaned, warnings = _sanitise_bullets(bullets)

    # Use a blank layout so we have full control
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    _apply_background(slide)

    _state["slide_count"] += 1
    slide_num = _state["slide_count"]

    theme = _state["theme"]

    # ── Title ──
    title_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(0.4), Inches(11.5), Inches(1.2),
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = TITLE_FONT_SIZE
    run.font.bold = True
    run.font.color.rgb = _hex_to_rgb(theme["title_color"])

    # ── Bullets ──
    body_box = slide.shapes.add_textbox(
        Inches(1.0), Inches(2.0), Inches(11.0), Inches(4.5),
    )
    tf_body = body_box.text_frame
    tf_body.word_wrap = True

    for idx, bullet in enumerate(cleaned):
        if idx == 0:
            p = tf_body.paragraphs[0]
        else:
            p = tf_body.add_paragraph()
        p.text = f"• {bullet}" if bullet else ""
        p.space_after = Pt(10)
        p.alignment = PP_ALIGN.LEFT
        if p.runs:
            run = p.runs[0]
            run.font.size = BODY_FONT_SIZE
            run.font.color.rgb = _hex_to_rgb(theme["body_color"])

    # ── Slide number footer ──
    _add_slide_number(slide, slide_num)

    logger.info("Slide %d added: %r (%d bullets).", slide_num, title, len(cleaned))
    return json.dumps({
        "status": "success",
        "slide_number": slide_num,
        "title": title,
        "bullet_count": len(cleaned),
        "warnings": warnings if warnings else None,
    })


@mcp.tool
def add_image_slide(title: str, image_path: str, caption: str = "") -> str:
    """
    Add a slide with a title, a centred image, and an optional caption.

    Args:
        title:      Slide title text.
        image_path: Absolute or relative path to the image file.
        caption:    Optional caption below the image.

    Returns:
        JSON string with slide details.
    """
    logger.info("add_image_slide(title=%r, image=%r)", title, image_path)

    err = _require_presentation()
    if err:
        return err

    if not os.path.isfile(image_path):
        logger.error("Image not found: %s", image_path)
        return json.dumps({
            "status": "error",
            "message": f"Image file not found: {image_path}",
        })

    prs: Presentation = _state["prs"]
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    _apply_background(slide)

    _state["slide_count"] += 1
    slide_num = _state["slide_count"]
    theme = _state["theme"]

    # ── Title ──
    title_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(0.3), Inches(11.5), Inches(1.0),
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = TITLE_FONT_SIZE
    run.font.bold = True
    run.font.color.rgb = _hex_to_rgb(theme["title_color"])

    # ── Image (centred) ──
    img_max_w = Inches(9.0)
    img_max_h = Inches(4.5)
    img_left = (WIDESCREEN_WIDTH - img_max_w) // 2
    img_top = Inches(1.5)
    slide.shapes.add_picture(image_path, img_left, img_top, img_max_w, img_max_h)

    # ── Caption ──
    if caption:
        cap_box = slide.shapes.add_textbox(
            Inches(1.0), Inches(6.2), Inches(11.0), Inches(0.6),
        )
        tf_cap = cap_box.text_frame
        tf_cap.word_wrap = True
        p_cap = tf_cap.paragraphs[0]
        p_cap.text = caption
        p_cap.alignment = PP_ALIGN.CENTER
        if p_cap.runs:
            run = p_cap.runs[0]
            run.font.size = Pt(16)
            run.font.italic = True
            run.font.color.rgb = _hex_to_rgb(theme["body_color"])

    # ── Slide number footer ──
    _add_slide_number(slide, slide_num)

    logger.info("Image slide %d added: %r", slide_num, title)
    return json.dumps({
        "status": "success",
        "slide_number": slide_num,
        "title": title,
        "image_path": image_path,
        "caption": caption or None,
    })


@mcp.tool
def set_theme(background_color: str, title_color: str, body_color: str) -> str:
    """
    Set the colour theme for subsequent slides.

    Colours are hex strings (e.g. 'FFFFFF', '#1A1A2E').
    This does NOT retroactively recolour existing slides.

    Args:
        background_color: Hex colour for slide backgrounds.
        title_color:      Hex colour for title text.
        body_color:       Hex colour for body / bullet text.

    Returns:
        JSON string confirming the new theme.
    """
    logger.info(
        "set_theme(bg=%s, title=%s, body=%s)",
        background_color, title_color, body_color,
    )

    # Validate all three colours
    for label, val in [
        ("background_color", background_color),
        ("title_color", title_color),
        ("body_color", body_color),
    ]:
        try:
            _hex_to_rgb(val)
        except ValueError as exc:
            logger.error("Invalid colour for %s: %s", label, exc)
            return json.dumps({
                "status": "error",
                "message": f"Invalid colour for {label}: {val}",
            })

    _state["theme"] = {
        "background_color": background_color.lstrip("#"),
        "title_color": title_color.lstrip("#"),
        "body_color": body_color.lstrip("#"),
    }

    logger.info("Theme updated.")
    return json.dumps({
        "status": "success",
        "message": "Theme updated. New slides will use these colours.",
        "theme": _state["theme"],
    })


@mcp.tool
def save_presentation(output_path: str) -> str:
    """
    Save the current presentation to disk as a .pptx file.

    Automatically creates parent directories if they don't exist.

    Args:
        output_path: File path for the saved .pptx.

    Returns:
        JSON string confirming save or describing the error.
    """
    logger.info("save_presentation(output_path=%r)", output_path)

    err = _require_presentation()
    if err:
        return err

    try:
        # Ensure output directory exists
        out_dir = os.path.dirname(os.path.abspath(output_path))
        os.makedirs(out_dir, exist_ok=True)

        _state["prs"].save(output_path)
        abs_path = os.path.abspath(output_path)
        logger.info("Presentation saved → %s (%d slides).", abs_path, _state["slide_count"])
        return json.dumps({
            "status": "success",
            "message": f"Presentation saved successfully.",
            "output_path": abs_path,
            "slide_count": _state["slide_count"],
        })
    except Exception:
        tb = traceback.format_exc()
        logger.error("Save failed:\n%s", tb)
        return json.dumps({
            "status": "error",
            "message": "Failed to save presentation.",
            "details": tb,
        })


# ──────────────────────────────────────────────
# Entry point
# ──────────────────────────────────────────────
if __name__ == "__main__":
    logger.info("Starting PPTX MCP Server …")
    mcp.run()
